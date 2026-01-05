from pptx import Presentation
from pptx.util import Inches, Pt
import os

BASE = os.path.dirname(os.path.abspath(__file__))
OUT = os.path.join(BASE, 'outputs')
PPT_PATH = os.path.join(OUT, 'Aadhaar_analytics_presentation.pptx')

prs = Presentation()

def add_title(title, subtitle=''):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    if subtitle:
        slide.placeholders[1].text = subtitle

def add_bullets(title, bullets):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for i, b in enumerate(bullets):
        if i == 0:
            p = body.paragraphs[0]
            p.text = b
        else:
            p = body.add_paragraph()
            p.text = b
        p.level = 0
        p.font.size = Pt(18)

def add_image_slide(title, image_path, left=Inches(0.5), top=Inches(1.6), width=Inches(9)):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    if os.path.exists(image_path):
        slide.shapes.add_picture(image_path, left, top, width=width)


def main():
    add_title('Identifying Regional and Age‑Group Patterns in Aadhaar Updates', 'Aadhaar Data Hackathon — Analytics')

    add_bullets('Objective', [
        'Identify regional (state/district) and age‑group patterns in Aadhaar updates',
        'Provide insights and recommendations to improve service delivery'
    ])

    add_bullets('Dataset & Method', [
        'Demographic and Biometric aggregated CSVs (daily rows)',
        'Aggregated by state/district; weekly resampling for forecasts',
        'Holt–Winters forecasts (12-week horizon) for key series'
    ])

    # Charts
    charts = [
        ('Top 10 States — Demographic', os.path.join(OUT, 'state_demographic_top10.png')),
        ('Top 10 States — Biometric', os.path.join(OUT, 'state_biometric_top10.png')),
        ('Gujarat — Districts (Demographic)', os.path.join(OUT, 'gujarat_demographic_top15.png')),
        ('Gujarat — Districts (Biometric)', os.path.join(OUT, 'gujarat_biometric_top15.png')),
    ]
    for title, path in charts:
        add_image_slide(title, path)

    # Forecast sample slides (top states)
    fc_dir = os.path.join(OUT, 'forecasts')
    sample_forecasts = [
        ('Forecast — Uttar Pradesh', os.path.join(fc_dir, 'state_Uttar_Pradesh_forecast.png')),
        ('Forecast — Maharashtra', os.path.join(fc_dir, 'state_Maharashtra_forecast.png')),
        ('Forecast — Gujarat / Ahmedabad', os.path.join(fc_dir, 'gujarat_Ahmedabad_forecast.png')),
    ]
    for title, path in sample_forecasts:
        add_image_slide(title, path)

    # Insights and recommendations (concise)
    add_bullets('Key Insights', [
        'Uttar Pradesh and Maharashtra show consistently high update volumes',
        'Ahmedabad and Surat dominate Gujarat update requests — urban concentration',
        'High biometric updates indicate revalidation/quality needs in populous areas'
    ])

    add_bullets('Recommendations', [
        'Scale temporary/mobile camps and staff during forecasted high weeks',
        'Targeted outreach in high-update Gujarat districts for address/mobile updates',
        'Implement biometric quality/revalidation programs to reduce repeat visits'
    ])

    add_bullets('Next Steps', [
        'Compare forecasts with ARIMA/Prophet models for accuracy',
        'Incorporate event calendars as covariates (campaigns, policy changes)',
        'Prepare PPT speaker notes and finalize slides for submission'
    ])

    # Final slide
    add_bullets('Conclusion', ['Data-driven actions can significantly improve UIDAI service delivery efficiency.'])

    prs.save(PPT_PATH)
    print('Presentation saved to', PPT_PATH)


if __name__ == '__main__':
    main()
